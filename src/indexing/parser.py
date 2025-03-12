import os
import logging
from typing import Dict, Optional
import fitz  # PyMuPDF for PDFs
import docx2txt  # For DOCX
from pptx import Presentation  # For PPTX
from concurrent.futures import ThreadPoolExecutor

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

BASE_DIR = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))  # src/indexing/ -> src/ -> root
DOCS_FOLDER = os.path.join(BASE_DIR, "data/docs")

def extract_from_pdf(file_path: str, use_parallel: bool = True, min_pages_for_parallel: int = 10) -> str:
    """Extract text from PDF using PyMuPDF (fastest)."""
    try:
        with fitz.open(file_path) as doc:
            if not doc.page_count:
                raise ValueError("PDF contains no pages.")
            if use_parallel and doc.page_count >= min_pages_for_parallel:
                # Use parallel processing for large PDFs
                with ThreadPoolExecutor() as executor:
                    results = executor.map(lambda page: page.get_text("text"), doc)
                text = list(results)
            else:
                # Use sequential processing for small PDFs
                text = [page.get_text("text") for page in doc]
            return "\n".join(text)
    except Exception as e:
        logger.error(f"Error extracting text from PDF {file_path}: {str(e)}")
        raise

def extract_from_docx(file_path: str) -> str:
    """Extract text from DOCX."""
    try:
        text = docx2txt.process(file_path)
        if not text.strip():
            raise ValueError("DOCX file contains no extractable text.")
        return text
    except Exception as e:
        logger.error(f"Error extracting text from DOCX {file_path}: {str(e)}")
        raise

def extract_from_pptx(file_path: str) -> str:
    """Extract text from PPTX slides."""
    try:
        text = []
        prs = Presentation(file_path)
        if not prs.slides:
            raise ValueError("PPTX file contains no slides.")
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        extracted_text = "\n".join(text)
        if not extracted_text.strip():
            raise ValueError("PPTX file contains no extractable text.")
        return extracted_text
    except Exception as e:
        logger.error(f"Error extracting text from PPTX {file_path}: {str(e)}")
        raise

def extract_text(file_path: str) -> Dict[str, Optional[str]]:
    """
    Detect file type and extract text accordingly.
    Returns a dictionary with 'text' and 'metadata' keys.
    """
    # Validate file path
    if not os.path.exists(file_path):
        logger.error(f"File does not exist: {file_path}")
        raise FileNotFoundError(f"File does not exist: {file_path}")
    if not os.path.isfile(file_path):
        logger.error(f"Path is not a file: {file_path}")
        raise ValueError(f"Path is not a file: {file_path}")

    # Get file metadata
    metadata = {
        "file_path": file_path,
        "file_name": os.path.basename(file_path),
        "file_size": os.path.getsize(file_path),
        "last_modified": os.path.getmtime(file_path),
    }

    # Detect file type
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    try:
        if ext == ".pdf":
            text = extract_from_pdf(file_path)
        elif ext == ".docx":
            text = extract_from_docx(file_path)
        elif ext == ".pptx":
            text = extract_from_pptx(file_path)
        else:
            logger.error(f"Unsupported file format: {file_path}")
            raise ValueError(f"Unsupported file format: {ext}")

        if not text.strip():
            logger.warning(f"No text extracted from file: {file_path}")
            raise ValueError(f"No text extracted from file: {file_path}")

        return {"text": text, "metadata": metadata}
    except Exception as e:
        logger.error(f"Failed to extract text from {file_path}: {str(e)}")
        raise

# Example Usage:
# if __name__ == "__main__":
#     sample_pdf = "data/docs/p2.pdf"
#     sample_docx = "data/docs/sample.docx"
#     sample_pptx = "data/docs/sample.pptx"

#     try:
#         pdf_result = extract_text(sample_pdf)
#         print("PDF Text:", pdf_result["text"][:500])  # Show first 500 chars
#         print("PDF Metadata:", pdf_result["metadata"])
#     except Exception as e:
#         print(f"PDF Error: {str(e)}")

#     try:
#         docx_result = extract_text(sample_docx)
#         print("DOCX Text:", docx_result["text"][:500])
#         print("DOCX Metadata:", docx_result["metadata"])
#     except Exception as e:
#         print(f"DOCX Error: {str(e)}")

#     try:
#         pptx_result = extract_text(sample_pptx)
#         print("PPTX Text:", pptx_result["text"][:500])
#         print("PPTX Metadata:", pptx_result["metadata"])
#     except Exception as e:
#         print(f"PPTX Error: {str(e)}")